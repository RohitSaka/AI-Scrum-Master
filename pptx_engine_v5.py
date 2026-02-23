"""
Sprint Deck Engine — Pixel-perfect replication of Bi-weekly Status Report.
Generates 7 slides: Hero, Agenda, Sprint Overview, KPIs, Accomplishments,
Risks & Blockers, Next Steps.

Data contract: generate_sprint_deck(data) where data is dict from Jira.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from io import BytesIO
import copy

# ════════════════════════════════════════════
#  CONSTANTS — exact from reference analysis
# ════════════════════════════════════════════
SW = 13.333  # slide width inches
SH = 7.500   # slide height inches
_i = lambda v: Inches(v)

# Colors
C_NAVY     = RGBColor(0x1E, 0x3A, 0x8A)  # header bars
C_BLUE     = RGBColor(0x3B, 0x82, 0xF6)  # accents
C_BLUE_D   = RGBColor(0x25, 0x63, 0xEB)  # POD blue
C_BLUE_HI  = RGBColor(0x1D, 0x4E, 0xD8)  # "Sprint Review" text
C_GREEN    = RGBColor(0x10, 0xB9, 0x81)  # POD green
C_PURPLE   = RGBColor(0x8B, 0x5C, 0xF6)  # POD purple
C_RED      = RGBColor(0xDC, 0x26, 0x26)  # risks red
C_RED_MIT  = RGBColor(0x05, 0x96, 0x69)  # mitigated green

C_BG       = RGBColor(0xF4, 0xF6, 0xF8)  # slide bg
C_CONTENT  = RGBColor(0xF9, 0xFA, 0xFB)  # content area
C_CARD     = RGBColor(0xFF, 0xFF, 0xFF)  # white cards
C_BORDER   = RGBColor(0xE5, 0xE7, 0xEB)  # card borders
C_GRAY_F3  = RGBColor(0xF3, 0xF4, 0xF6)  # stats bar bg

C_TXT_DARK = RGBColor(0x1F, 0x29, 0x37)  # primary text
C_TXT_BLK  = RGBColor(0x11, 0x18, 0x27)  # title text
C_TXT_MID  = RGBColor(0x4B, 0x55, 0x63)  # secondary text
C_TXT_GRAY = RGBColor(0x6B, 0x72, 0x80)  # muted text
C_TXT_LT   = RGBColor(0x9C, 0xA3, 0xAF)  # footer text
C_GHOST    = RGBColor(0xE5, 0xE7, 0xEB)  # ghost numbers
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# POD accent backgrounds
C_BLUE_BG  = RGBColor(0xDB, 0xEA, 0xFE)  # blue light bg
C_BLUE_BG2 = RGBColor(0xEF, 0xF6, 0xFF)  # blue pill bg
C_GREEN_BG = RGBColor(0xD1, 0xFA, 0xE5)  # green light bg
C_PURP_BG  = RGBColor(0xED, 0xE9, 0xFE)  # purple light bg
C_RED_BG   = RGBColor(0xFE, 0xCA, 0xCA)  # risk red bg
C_GREEN_MIT= RGBColor(0xA7, 0xF3, 0xD0)  # mitigated green bg

C_BLUE_TXT = RGBColor(0x1E, 0x40, 0xAF)  # POD badge text
C_PURP_TXT = RGBColor(0x7C, 0x3A, 0xED)  # purple badge text
C_TXT_374  = RGBColor(0x37, 0x41, 0x51)  # bullet text

# Layout
HEADER_H   = 0.833   # navy header bar height
FOOTER_H   = 0.510   # footer height
FOOTER_Y   = SH - FOOTER_H  # 6.990
MX         = 0.625   # horizontal margin
COL_W      = 4.031   # POD column width
COL_GAP    = 0.236   # gap between columns
COL_X      = [0.417, 4.653, 8.889]  # column left positions
CARD_Y     = 1.146   # card top
CARD_H     = 5.531   # card height (overview/accomplishments/next)
CARD_H2    = 6.229   # card height (KPI slide - extends further)

# Fonts
F_HEAD = 'Montserrat'
F_BODY = 'Roboto'


# ════════════════════════════════════════════
#  HELPERS
# ════════════════════════════════════════════
def _rect(s, x, y, w, h, fill=None, border=None, border_w=1, radius=0):
    """Add rectangle shape."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                            _i(x), _i(y), _i(w), _i(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if border:
        sh.line.color.rgb = border; sh.line.width = Pt(border_w)
    else:
        sh.line.fill.background()
    if radius:
        try:
            sh._element.attrib['{http://schemas.microsoft.com/office/drawing/2010/main}' + 'adj'] = str(radius)
        except:
            pass
    return sh

def _txt(s, text, x, y, w, h, sz, color, bold=False, font=F_BODY,
         align=PP_ALIGN.LEFT, v_anchor=None, wrap=True):
    """Add text box."""
    if not text:
        return None
    tb = s.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    if v_anchor:
        try:
            tf.vertical_anchor = v_anchor
        except:
            pass
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb

def _txt_multi(s, text, x, y, w, h, sz, color, bold=False, font=F_BODY, align=PP_ALIGN.LEFT):
    """Add text box with auto line wrapping."""
    tb = s.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb

def _line_h(s, x, y, w, h, color):
    """Horizontal divider line."""
    return _rect(s, x, y, w, max(h, 0.01), fill=color, border=color)

def _line_v(s, x, y, w, h, color):
    """Vertical accent bar."""
    return _rect(s, x, y, max(w, 0.02), h, fill=color, border=color)

def _labeled_shape(s, x, y, w, h, fill, text, sz, text_color, font=F_HEAD, bold=True):
    """Rectangle with centered text (for badges, buttons, headers)."""
    sh = _rect(s, x, y, w, h, fill=fill)
    sh.line.fill.background()
    tb = _txt(s, text, x, y, w, h, sz, text_color, bold=bold, font=font,
              align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    return sh

def _add_table(s, x, y, w, rows_data, col_widths_pct, header_color=C_BLUE_D):
    """Add a formatted table. rows_data = list of lists. First row = header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 3
    row_h = Inches(0.28)

    tbl_shape = s.shapes.add_table(n_rows, n_cols, _i(x), _i(y),
                                    _i(w), row_h * n_rows)
    tbl = tbl_shape.table

    # Set column widths
    total_w = w
    for ci in range(n_cols):
        tbl.columns[ci].width = _i(total_w * col_widths_pct[ci])

    for ri, row_data in enumerate(rows_data):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)

            # Style
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.name = F_BODY
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = C_WHITE
                    p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                else:
                    p.font.bold = ci == 0
                    p.font.color.rgb = C_TXT_DARK if ci == 0 else C_TXT_MID
                    p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT

            # Cell fill
            cf = cell.fill
            if ri == 0:
                cf.solid(); cf.fore_color.rgb = header_color
            elif ri % 2 == 0:
                cf.solid(); cf.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
            else:
                cf.solid(); cf.fore_color.rgb = C_WHITE

            cell.margin_left = Emu(45720)
            cell.margin_right = Emu(45720)
            cell.margin_top = Emu(18288)
            cell.margin_bottom = Emu(18288)

    # Remove table borders for clean look
    tbl_xml = tbl._tbl
    for tc in tbl_xml.iter(qn('a:tc')):
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is None:
            tcPr = tc.makeelement(qn('a:tcPr'), {})
            tc.insert(0, tcPr)
        for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
            ln = tcPr.find(qn(f'a:{border_name}'))
            if ln is None:
                ln = tcPr.makeelement(qn(f'a:{border_name}'), {'w': '6350'})
                tcPr.append(ln)
            else:
                ln.set('w', '6350')
            sf = ln.find(qn('a:solidFill'))
            if sf is None:
                sf = ln.makeelement(qn('a:solidFill'), {})
                ln.append(sf)
            srgb = sf.find(qn('a:srgbClr'))
            if srgb is None:
                srgb = sf.makeelement(qn('a:srgbClr'), {'val': 'E5E7EB'})
                sf.append(srgb)
            else:
                srgb.set('val', 'E5E7EB')

    return tbl_shape


# ════════════════════════════════════════════
#  COMMON SLIDE ELEMENTS
# ════════════════════════════════════════════
def _slide_bg(s):
    """Full slide background #F4F6F8."""
    _rect(s, 0, 0, SW, SH, fill=C_BG, border=C_WHITE)

def _header_bar(s, title):
    """Navy header bar with white title."""
    _rect(s, 0, 0, SW, HEADER_H, fill=C_NAVY, border=C_WHITE)
    _txt(s, title, MX, 0.167, 6, 0.5, 24, C_WHITE, bold=True, font=F_HEAD)

def _content_bg(s, h=None):
    """Light content background below header."""
    ch = h or (FOOTER_Y - HEADER_H)
    _rect(s, 0, HEADER_H, SW, ch, fill=C_CONTENT, border=C_WHITE)

def _footer(s, report_name, date_range):
    """Footer with report name left, date right."""
    _rect(s, 0, FOOTER_Y, SW, FOOTER_H, fill=C_WHITE)
    _line_h(s, 0, FOOTER_Y, SW, 0.01, C_BORDER)
    _txt(s, report_name, MX, FOOTER_Y + 0.166, 3, 0.188, 9, C_TXT_LT, font=F_BODY)
    _txt(s, date_range, 11.234, FOOTER_Y + 0.166, 1.761, 0.188, 9, C_TXT_LT, font=F_BODY)

def _pod_card(s, col_idx, card_h=CARD_H):
    """White card for a POD column."""
    x = COL_X[col_idx]
    _rect(s, x, CARD_Y, COL_W, card_h, fill=C_CARD, border=C_BORDER)

def _pod_header(s, col_idx, text, color, h=0.625):
    """Colored POD header band at top of card."""
    x = COL_X[col_idx] + 0.01
    _labeled_shape(s, x, CARD_Y + 0.01, COL_W - 0.02, h, color, text, 13, C_WHITE, font=F_HEAD)

def _section_label(s, x, y, text):
    """Gray section label (e.g. 'Team Roles', 'Sprint Focus')."""
    sh = _rect(s, x, y, 3.719, 0.188, fill=C_WHITE)
    sh.line.fill.background()
    _txt(s, text, x, y, 3.719, 0.188, 9, C_TXT_LT, bold=True, font=F_HEAD)

def _section_divider(s, x, y, w=3.594):
    """Thin gray horizontal divider."""
    _line_h(s, x, y, w, 0.01, C_BORDER)


# ════════════════════════════════════════════
#  POD DATA HELPERS
# ════════════════════════════════════════════
POD_COLORS = [C_BLUE_D, C_GREEN, C_PURPLE]
POD_BG_COLORS = [C_BLUE_BG, C_GREEN_BG, C_PURP_BG]
POD_TXT_COLORS = [C_BLUE_D, C_GREEN, C_PURP_TXT]

def _get_pods(data):
    """Get pods list, pad to 3 if needed."""
    pods = data.get('pods', [])
    while len(pods) < 3:
        pods.append({'name': f'POD {len(pods)+1}', 'short_name': f'POD {len(pods)+1}',
                     'sprint_focus': [], 'accomplishments': [], 'risks': [], 'next_steps': [],
                     'total_issues': 0, 'story_points': 0, 'issue_breakdown': [], 'story_points_alloc': []})
    return pods[:3]


# ════════════════════════════════════════════
#  SLIDE 1: HERO / COVER
# ════════════════════════════════════════════
def build_hero(s, data):
    """Cover slide — split layout with diagonal blue shapes."""
    _slide_bg(s)

    # Right diagonal shape 1 — lighter blue (#3B82F6)
    # Freeform diagonal from reference: starts at ~6.93" from left
    _rect(s, 6.93, 0, 6.5, SH, fill=C_BLUE)

    # Right diagonal shape 2 — dark navy (#1E3A8A)
    _rect(s, 7.33, 0, 6.1, SH, fill=C_NAVY)

    # "Sprint Review" eyebrow label
    eyebrow = data.get('eyebrow', 'Sprint Review')
    _txt(s, eyebrow, 0.833, 1.220, 7.292, 0.3, 10, C_BLUE_HI, bold=True, font=F_BODY)

    # Main title — two lines
    line1 = data.get('title_line1', 'BI-Weekly Project')
    line2 = data.get('title_line2', 'Status Report')
    _txt(s, line1, 0.833, 1.595, 7.355, 0.7, 48, C_TXT_BLK, bold=True, font=F_HEAD)
    _txt(s, line2, 0.833, 2.3, 7.355, 0.7, 48, C_TXT_BLK, bold=True, font=F_HEAD)

    # Left vertical accent bar
    _line_v(s, 0.833, 3.312, 0.063, 1.313, C_NAVY)

    # Description text
    desc = data.get('subtitle', 'Comprehensive overview of sprint progress, metrics')
    _txt_multi(s, desc, 1.104, 3.312, 5.5, 1.313, 21, C_TXT_MID, font=F_HEAD)

    # Calendar icon + date
    _labeled_shape(s, 0.833, 5.140, 0.250, 0.250, C_BORDER, '\U0001F4C5', 9, C_NAVY, font=F_BODY)
    sprint = data.get('sprint_name', 'FY26.PI3.S2')
    date_range = data.get('date_range', '02/05/2026 – 02/18/2026')
    # Two-part date text: bold period + regular dates
    tb = s.shapes.add_textbox(_i(1.208), _i(5.124), _i(4.993), _i(0.282))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f'Reporting Period: {sprint}'
    run1.font.size = Pt(13); run1.font.color.rgb = C_TXT_GRAY
    run1.font.bold = True; run1.font.name = F_BODY
    run2 = p.add_run()
    run2.text = f'  [{date_range}]'
    run2.font.size = Pt(13); run2.font.color.rgb = C_TXT_GRAY
    run2.font.bold = False; run2.font.name = F_BODY

    # POD badges at bottom
    pods = _get_pods(data)
    badge_y = 5.843
    badge_h = 0.438
    badge_x = 0.833
    for i, pod in enumerate(pods):
        name = pod.get('name', f'POD {i+1}')
        bw = max(len(name) * 0.12 + 0.4, 1.0)
        sh = _rect(s, badge_x, badge_y, bw, badge_h, fill=C_BLUE_BG2, border=C_BLUE_BG)
        _txt(s, name, badge_x, badge_y, bw, badge_h, 12, C_BLUE_TXT, bold=True,
             font=F_BODY, align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
        badge_x += bw + 0.12

    # Company name bottom-right
    company = data.get('company', 'GENSPARK CONSULTING')
    _txt(s, company, 10.554, 6.729, 2.243, 0.25, 12, C_NAVY, bold=True, font=F_HEAD)

    # Bottom accent line
    _line_h(s, 0, SH - 0.04, SW, 0.04, C_BLUE)


# ════════════════════════════════════════════
#  SLIDE 2: AGENDA
# ════════════════════════════════════════════
AGENDA_ITEMS = [
    ('Sprint Overview', 'Team structure and current focus areas'),
    ('KPIs & Story Count', 'Issue breakdown and velocity metrics'),
    ('Accomplishments', 'Key achievements and delivered value'),
    ('Risks & Blockers', 'Current challenges and mitigation plans'),
    ('Next Steps', 'Upcoming priorities for the next sprint'),
]

def build_agenda(s, data):
    """Agenda slide — numbered cards in 2-column grid."""
    _slide_bg(s)

    # Navy header
    _rect(s, 0, 0, SW, 1.042, fill=C_NAVY, border=C_WHITE)
    _txt(s, 'Agenda', MX, 0.240, 2, 0.563, 27, C_WHITE, bold=True, font=F_HEAD)

    # Footer
    report_name = data.get('report_name', 'Weekly Project Status Report')
    date_range = data.get('date_range', '02/05/2026 – 02/18/2026')
    _footer(s, report_name, date_range)

    # Agenda cards — 2 columns, 3 rows (5 items)
    items = data.get('agenda_items', AGENDA_ITEMS)
    card_w = 5.781
    card_h = 1.083
    gap_x = 6.302  # col2 start - col1 start
    gap_y = 1.344   # vertical gap between rows
    start_x = MX
    start_y = 2.130

    for i, (title, desc) in enumerate(items):
        row = i // 2
        col = i % 2
        cx = start_x + col * gap_x
        cy = start_y + row * (card_h + gap_y - card_h + 0.229)  # row spacing

        # Card background
        _rect(s, cx, cy, card_w, card_h, fill=C_WHITE)

        # Left blue accent bar
        _rect(s, cx, cy, 0.052, card_h, fill=C_BLUE, border=C_BLUE)

        # Ghost number (right-aligned within left zone)
        num = f'{i+1:02d}'
        _txt(s, num, cx + 0.177, cy + 0.267, 0.719, 0.5, 36, C_GHOST,
             bold=True, font=F_HEAD, align=PP_ALIGN.RIGHT)

        # Title
        _txt(s, title, cx + 1.094, cy + 0.209, 3.5, 0.375, 18, C_TXT_DARK,
             bold=True, font=F_HEAD)

        # Description
        _txt(s, desc, cx + 1.094, cy + 0.625, 3.5, 0.25, 12, C_TXT_GRAY, font=F_BODY)


# ════════════════════════════════════════════
#  SLIDE 3: SPRINT OVERVIEW
# ════════════════════════════════════════════
def build_sprint_overview(s, data):
    """Sprint Overview — 3 POD columns with team roles + sprint focus."""
    _slide_bg(s)
    _header_bar(s, 'Sprint Overview')
    _content_bg(s)

    report_name = data.get('report_name', 'Weekly Project Status Report')
    date_range = data.get('date_range', '02/05/2026 – 02/18/2026')
    _footer(s, report_name, date_range)

    pods = _get_pods(data)
    for ci, pod in enumerate(pods):
        x = COL_X[ci]
        color = POD_COLORS[ci]
        inner_x = x + 0.218  # inner content margin

        # Card + colored header
        _pod_card(s, ci)
        _pod_header(s, ci, pod.get('name', f'POD {ci+1}'), color)

        # TEAM ROLES section
        _section_label(s, inner_x, 1.990, 'Team Roles')
        _section_divider(s, inner_x, 2.230)

        roles = [
            ('Scrum Mgr:', pod.get('scrum_mgr', '—')),
            ('SEM:', pod.get('sem', '—')),
            ('Product Mgr:', pod.get('product_mgr', '—')),
        ]
        role_y = 2.344
        for label, value in roles:
            # Role icon placeholder (small bullet)
            _txt(s, '●', inner_x, role_y, 0.15, 0.209, 7, C_TXT_MID, font=F_BODY)
            _txt(s, label, inner_x + 0.22, role_y, 1.23, 0.209, 9, C_TXT_MID, bold=True, font=F_BODY)
            _txt(s, value, inner_x + 1.47, role_y, 1.5, 0.209, 9, C_TXT_DARK, font=F_BODY)
            role_y += 0.286

        # SPRINT FOCUS section
        _section_label(s, inner_x, 3.411, 'Sprint Focus')
        _section_divider(s, inner_x, 3.651)

        focus_items = pod.get('sprint_focus', [])
        fy = 3.766
        for item in focus_items[:6]:
            item_text = item if isinstance(item, str) else str(item.get('title', item))
            _txt(s, '●', inner_x, fy + 0.04, 0.1, 0.15, 6, C_BLUE_D, font=F_BODY)
            _txt(s, item_text, inner_x + 0.19, fy, 3.4, 0.209, 10, C_TXT_374, font=F_BODY)
            fy += 0.32


# ════════════════════════════════════════════
#  SLIDE 4: KPIs & STORY COUNT
# ════════════════════════════════════════════
def build_kpis(s, data):
    """KPIs & Story Count — 3 POD columns with metrics + tables."""
    _slide_bg(s)
    _header_bar(s, 'KPIs & Story Count')
    _content_bg(s, h=SH - HEADER_H)

    report_name = data.get('report_name', 'Weekly Project Status Report')
    date_range = data.get('date_range', '02/05/2026 – 02/18/2026')
    # Footer at bottom (extended slide)
    _rect(s, 0, SH - 0.04, SW, 0.04, fill=C_BLUE)

    pods = _get_pods(data)
    for ci, pod in enumerate(pods):
        x = COL_X[ci]
        color = POD_COLORS[ci]
        inner_x = x + 0.166

        # Card + header (shorter header for KPIs)
        _rect(s, x, 1.094, COL_W, CARD_H2, fill=C_CARD, border=C_BORDER)
        short_name = pod.get('short_name', pod.get('name', f'POD {ci+1}'))
        _labeled_shape(s, x + 0.01, 1.104, COL_W - 0.02, 0.469,
                       color, short_name, 12, C_WHITE, font=F_HEAD)

        # Stats bar (gray bg with Total Issues + Story Points)
        stats_y = 1.729
        _rect(s, inner_x, stats_y, 3.698, 0.656, fill=C_GRAY_F3, border=C_WHITE)

        total_issues = pod.get('total_issues', 0)
        story_points = pod.get('story_points', 0)

        # Total Issues (left half)
        _txt(s, str(total_issues), inner_x + 0.584, stats_y + 0.104, 0.771, 0.25, 18,
             color, bold=True, font=F_HEAD, align=PP_ALIGN.CENTER)
        _txt(s, 'Total Issues', inner_x + 0.580, stats_y + 0.396, 0.778, 0.157, 8,
             C_TXT_GRAY, font=F_BODY, align=PP_ALIGN.CENTER)

        # Story Points (right half)
        _txt(s, str(story_points), inner_x + 2.11, stats_y + 0.104, 0.803, 0.25, 18,
             color, bold=True, font=F_HEAD, align=PP_ALIGN.CENTER)
        _txt(s, 'Story Points', inner_x + 2.11, stats_y + 0.396, 0.803, 0.157, 8,
             C_TXT_GRAY, font=F_BODY, align=PP_ALIGN.CENTER)

        # Issue Breakdown section
        _txt(s, '■', inner_x, 2.635, 0.15, 0.15, 7, color, font=F_BODY)
        _txt(s, ' Issue Breakdown ', inner_x + 0.18, 2.594, 3.5, 0.188, 9,
             C_TXT_374, bold=True, font=F_HEAD)

        # Issue breakdown table
        breakdown = pod.get('issue_breakdown', [])
        if breakdown:
            rows = [['Type', 'To Do', 'In Progress', 'Done', 'Total']]
            for item in breakdown:
                rows.append([
                    item.get('type', ''),
                    str(item.get('todo', 0)),
                    str(item.get('progress', 0)),
                    str(item.get('done', 0)),
                    str(item.get('total', 0))
                ])
            _add_table(s, inner_x, 2.865, 3.694, rows, [0.35, 0.16, 0.19, 0.14, 0.16],
                       header_color=color)

        # Story Points Allocation section
        _txt(s, '■', inner_x, 4.693, 0.15, 0.15, 7, color, font=F_BODY)
        _txt(s, ' Story Points Allocation ', inner_x + 0.18, 4.651, 3.5, 0.188, 9,
             C_TXT_374, bold=True, font=F_HEAD)

        sp_alloc = pod.get('story_points_alloc', [])
        if sp_alloc:
            rows2 = [['Epic / Feature', 'Points', 'Status']]
            for item in sp_alloc:
                rows2.append([
                    item.get('epic', ''),
                    str(item.get('points', 0)),
                    item.get('status', '—')
                ])
            _add_table(s, inner_x, 4.922, 3.695, rows2, [0.50, 0.22, 0.28],
                       header_color=color)


# ════════════════════════════════════════════
#  SLIDE 5: ACCOMPLISHMENTS
# ════════════════════════════════════════════
def build_accomplishments(s, data):
    """Accomplishments — 3 POD columns with achievement items."""
    _slide_bg(s)
    _header_bar(s, 'Accomplishments')
    _content_bg(s)

    report_name = data.get('report_name', 'Weekly Project Status Report')
    date_range = data.get('date_range', '02/05/2026 – 02/18/2026')
    _footer(s, report_name, date_range)

    pods = _get_pods(data)
    for ci, pod in enumerate(pods):
        x = COL_X[ci]
        color = POD_COLORS[ci]
        bg_color = POD_BG_COLORS[ci]
        inner_x = x + 0.218

        # Card + header
        _pod_card(s, ci)
        short_name = pod.get('short_name', pod.get('name', f'POD {ci+1}'))
        _pod_header(s, ci, short_name, color)

        # Accomplishment items
        items = pod.get('accomplishments', [])
        iy = 2.063
        for j, item in enumerate(items[:6]):
            text = item if isinstance(item, str) else str(item.get('text', item))

            # Colored icon badge
            _rect(s, inner_x, iy, 0.25, 0.25, fill=bg_color, border=C_WHITE)
            _txt(s, '✓', inner_x + 0.05, iy + 0.02, 0.15, 0.21, 9, color,
                 bold=True, font=F_BODY, align=PP_ALIGN.CENTER)

            # Text
            _txt_multi(s, text, inner_x + 0.4, iy - 0.03, 3.271, 0.65, 10, C_TXT_374, font=F_BODY)

            # Spacing depends on text length
            lines_est = max(1, len(text) // 45 + 1)
            iy += 0.25 + lines_est * 0.22 + 0.15


# ════════════════════════════════════════════
#  SLIDE 6: RISKS & BLOCKERS
# ════════════════════════════════════════════
def build_risks(s, data):
    """Risks & Blockers — split: Active (red, left) + Mitigated (green, right)."""
    _slide_bg(s)
    _header_bar(s, 'Risks & Blockers')
    _content_bg(s)

    report_name = data.get('report_name', 'Weekly Project Status Report')
    date_range = data.get('date_range', '02/05/2026 – 02/18/2026')
    _footer(s, report_name, date_range)

    half_w = 5.833
    left_x = MX
    right_x = 6.875

    # ── LEFT: Active Attention Items ──
    _txt(s, '⚠', left_x, 1.198, 0.2, 0.208, 10, C_RED, font=F_BODY)
    _txt(s, 'Active Attention Items', left_x + 0.313, 1.146, 3.5, 0.49, 15,
         C_RED, bold=True, font=F_HEAD)
    _line_h(s, left_x, 1.615, half_w, 0.021, C_RED_BG)

    # ── RIGHT: Mitigated & Watchlist ──
    _txt(s, '✓', right_x, 1.198, 0.2, 0.208, 10, C_RED_MIT, font=F_BODY)
    _txt(s, 'Mitigated & Watchlist', right_x + 0.333, 1.146, 3.5, 0.49, 15,
         C_RED_MIT, bold=True, font=F_HEAD)
    _line_h(s, right_x, 1.615, half_w, 0.021, C_GREEN_MIT)

    # Collect risks by status
    pods = _get_pods(data)
    active_risks = []
    mitigated_risks = []

    for ci, pod in enumerate(pods):
        color = POD_COLORS[ci]
        bg_color = POD_BG_COLORS[ci]
        txt_color = POD_TXT_COLORS[ci]
        pod_name = pod.get('name', f'POD {ci+1}')
        for risk in pod.get('risks', []):
            entry = {**risk, 'pod_name': pod_name, 'pod_color': color,
                     'pod_bg': bg_color, 'pod_txt': txt_color}
            status = risk.get('status', 'Active Risk').lower()
            if 'mitigat' in status or 'watch' in status or 'resolved' in status:
                mitigated_risks.append(entry)
            else:
                active_risks.append(entry)

    def _draw_risk_card(s, risk, base_x, y, half_w):
        """Draw a single risk card."""
        card_x = base_x + 0.260
        # White card background
        _rect(s, base_x, y, half_w, 1.85, fill=C_CARD, border=C_BORDER, border_w=0.5)

        # POD badge
        pod_name = risk['pod_name']
        bw = max(len(pod_name) * 0.09 + 0.3, 0.688)
        _labeled_shape(s, card_x, y + 0.208, bw, 0.260,
                       risk['pod_bg'], pod_name, 8, risk['pod_txt'], font=F_BODY)

        # Title
        _txt(s, risk.get('title', ''), card_x, y + 0.641, half_w - 0.6, 0.24, 12,
             C_TXT_DARK, bold=True, font=F_HEAD)

        # Description
        _txt_multi(s, risk.get('description', ''), card_x, y + 0.999, half_w - 0.6, 0.5,
                   10, C_TXT_MID, font=F_BODY)

        # Status badge
        status = risk.get('status', 'Active Risk')
        is_active = 'active' in status.lower() or 'block' in status.lower()
        badge_color = C_RED if is_active else C_RED_MIT
        badge_border = C_RED_BG if is_active else C_GREEN_MIT

        badge_w = max(len(status) * 0.07 + 0.3, 1.1)
        _rect(s, card_x + 0.29, y + 1.56, badge_w, 0.21, fill=C_WHITE, border=badge_border)
        _txt(s, f' {status} ', card_x + 0.29, y + 1.56, badge_w, 0.21, 9,
             badge_color, bold=True, font=F_BODY, align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    # Draw active risks
    ry = 1.948
    for risk in active_risks[:3]:
        _draw_risk_card(s, risk, left_x, ry, half_w)
        ry += 2.05

    # Draw mitigated risks
    ry = 1.948
    for risk in mitigated_risks[:3]:
        _draw_risk_card(s, risk, right_x, ry, half_w)
        ry += 2.05


# ════════════════════════════════════════════
#  SLIDE 7: NEXT STEPS
# ════════════════════════════════════════════
def build_next_steps(s, data):
    """Next Steps — 3 POD columns with task items."""
    _slide_bg(s)
    _header_bar(s, 'Next Steps')
    _content_bg(s)

    report_name = data.get('report_name', 'Weekly Project Status Report')
    date_range = data.get('date_range', '02/05/2026 – 02/18/2026')
    _footer(s, report_name, date_range)

    pods = _get_pods(data)
    for ci, pod in enumerate(pods):
        x = COL_X[ci]
        color = POD_COLORS[ci]
        inner_x = x + 0.218

        # Card + header
        _pod_card(s, ci)
        short_name = pod.get('short_name', pod.get('name', f'POD {ci+1}'))
        _pod_header(s, ci, short_name, color, h=0.573)

        # Task items
        items = pod.get('next_steps', [])
        iy = 1.990
        for j, item in enumerate(items[:4]):
            if isinstance(item, str):
                title = item
                desc = ''
            else:
                title = str(item.get('title', ''))
                desc = str(item.get('description', ''))

            # Colored circle icon
            _rect(s, inner_x, iy, 0.333, 0.333, fill=color, border=None, radius=50)
            _txt(s, '→', inner_x + 0.08, iy + 0.05, 0.17, 0.23, 12, C_WHITE,
                 bold=True, font=F_BODY, align=PP_ALIGN.CENTER)

            # Vertical accent line
            _line_v(s, inner_x + 0.167, iy + 0.365, 0.021, 0.292, color)

            # Title
            _txt(s, title, inner_x + 0.5, iy + 0.04, 3.196, 0.219, 10,
                 C_TXT_DARK, bold=True, font=F_HEAD)

            # Description
            if desc:
                _txt_multi(s, desc, inner_x + 0.5, iy + 0.292, 3.157, 0.407, 9,
                           C_TXT_MID, font=F_BODY)

            iy += 0.95


# ════════════════════════════════════════════
#  MAIN GENERATOR
# ════════════════════════════════════════════
SLIDE_BUILDERS = [
    build_hero,
    build_agenda,
    build_sprint_overview,
    build_kpis,
    build_accomplishments,
    build_risks,
    build_next_steps,
]

def generate_sprint_deck(data):
    """
    Generate a complete Sprint Review deck from Jira data.

    Args:
        data: dict with keys:
            sprint_name, date_range, company, title_line1, title_line2, subtitle,
            report_name, eyebrow,
            pods: list of pod dicts with:
                name, short_name, color,
                scrum_mgr, sem, product_mgr,
                sprint_focus: list of str,
                total_issues: int, story_points: int,
                issue_breakdown: list of {type, todo, progress, done, total},
                story_points_alloc: list of {epic, points, status},
                accomplishments: list of str or {text},
                risks: list of {title, description, status},
                next_steps: list of {title, description}

    Returns:
        BytesIO buffer with .pptx content
    """
    prs = Presentation()
    prs.slide_width = _i(SW)
    prs.slide_height = _i(SH)

    blank_layout = prs.slide_layouts[6]  # blank

    for builder in SLIDE_BUILDERS:
        slide = prs.slides.add_slide(blank_layout)
        builder(slide, data)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ════════════════════════════════════════════
#  SAMPLE DATA (for testing)
# ════════════════════════════════════════════
SAMPLE_DATA = {
    "sprint_name": "FY26.PI3.S2",
    "date_range": "02/05/2026 – 02/18/2026",
    "company": "GENSPARK CONSULTING",
    "title_line1": "BI-Weekly Project",
    "title_line2": "Status Report",
    "eyebrow": "Sprint Review",
    "subtitle": "Comprehensive overview of sprint progress, metrics and team activities across all active PODs",
    "report_name": "Weekly Project Status Report",
    "pods": [
        {
            "name": "Business Development POD",
            "short_name": "Business Development",
            "scrum_mgr": "Neetha Kotagiri",
            "sem": "Mothilal Ramavath",
            "product_mgr": "Lance Nazer",
            "sprint_focus": [
                "Transfer Center",
                "AscOne PEI batch scaling",
                "PRM NPI Search",
                "PRM Enhancements and Cleanup"
            ],
            "total_issues": 13,
            "story_points": 36,
            "issue_breakdown": [
                {"type": "Story", "todo": 0, "progress": 5, "done": 3, "total": 8},
                {"type": "Task", "todo": 1, "progress": 0, "done": 2, "total": 3},
                {"type": "Bug", "todo": 0, "progress": 1, "done": 1, "total": 2}
            ],
            "story_points_alloc": [
                {"epic": "Transfer Center", "points": 13, "status": "In Progress"},
                {"epic": "PEI Batch Scaling", "points": 8, "status": "In Progress"},
                {"epic": "PRM NPI Search", "points": 8, "status": "To Do"},
                {"epic": "PRM Enhancements", "points": 7, "status": "Done"}
            ],
            "accomplishments": [
                "Transfer center Enhancements: Implemented a standardized call flow and improved routing logic",
                "AscOne PEI batch scaling: Completed the initial batch processing framework",
                "PRM NPI Search: Built the NPI validation search module",
                "PRM Cleanup: Resolved 7 legacy data quality issues"
            ],
            "risks": [
                {"title": "PEI Scaling: Implementation for SFMC",
                 "description": "Blocked as we are waiting on the approval for the design from GCP team. Will resume implementation once approved.",
                 "status": "Active Risk"},
                {"title": "GCP Team Design Issue",
                 "description": "Design issue identified with the Google Cloud Platform integration. Team has proposed alternative architecture.",
                 "status": "Mitigated"}
            ],
            "next_steps": [
                {"title": "BD – Gmail Plug In", "description": "Discovery for identifying right plug in for generating leads from email interactions"},
                {"title": "ACVS – Solution discovery", "description": "Document Discovery for ACVS implementation on salesforce platform"},
                {"title": "AI Summaries Discovery", "description": "Initiate discovery for AI-driven tasks and cases summarization capabilities"}
            ]
        },
        {
            "name": "Provider Services POD",
            "short_name": "Provider Services",
            "scrum_mgr": "Neetha Kotagiri",
            "sem": "Surya Chinta",
            "product_mgr": "Nichole Van",
            "sprint_focus": [
                "Form Assembly: Provider Type logic & general info",
                "Form Assembly: Billing/Remittance/Identification",
                "Texas CIN migration: solution design",
                "Automated reminder email enhancements",
                "Legacy report migration to Enterprise"
            ],
            "total_issues": 16,
            "story_points": 49,
            "issue_breakdown": [
                {"type": "Story", "todo": 2, "progress": 6, "done": 4, "total": 12},
                {"type": "Task", "todo": 0, "progress": 1, "done": 2, "total": 3},
                {"type": "Bug", "todo": 0, "progress": 0, "done": 1, "total": 1}
            ],
            "story_points_alloc": [
                {"epic": "Form Assembly", "points": 18, "status": "In Progress"},
                {"epic": "Texas CIN Migration", "points": 13, "status": "In Progress"},
                {"epic": "Reminder Emails", "points": 10, "status": "In Progress"},
                {"epic": "Legacy Reports", "points": 5, "status": "To Do"},
                {"epic": "Provider Search", "points": 3, "status": "Done"}
            ],
            "accomplishments": [
                "Enabled provider-type driven dynamic intake (Professional vs Facility) with conditional field rendering",
                "Completed Texas CIN migration solution design document",
                "Automated reminder email enhancements — added configurable frequency settings",
                "Legacy report migration: Mapped 12 reports to Enterprise reporting framework",
                "Address APDH team request — Successfully sending Provider data via API integration"
            ],
            "risks": [],
            "next_steps": [
                {"title": "New Sharing Structure", "description": "Configure the updated sharing set-up to support Texas provider groups"},
                {"title": "Digital Provider App", "description": "Develop new sections for Licensure, Compliance, and Credentialing"},
                {"title": "Automated Reminder Emails", "description": "Complete the Automated email configuration to enable scheduled provider notifications"}
            ]
        },
        {
            "name": "VBC POD",
            "short_name": "VBC POD",
            "scrum_mgr": "Sandeep Kandiraju",
            "sem": "Suraj Chalanti",
            "product_mgr": "Kammal Sunmola",
            "sprint_focus": [
                "Ingest targeting data from ADSI into Data Cloud",
                "Setup ADSI data streaming and mapping to objects",
                "Segment members using ADSI indicators (CHF Tags)"
            ],
            "total_issues": 5,
            "story_points": 29,
            "issue_breakdown": [
                {"type": "Story", "todo": 0, "progress": 3, "done": 1, "total": 4},
                {"type": "Task", "todo": 0, "progress": 0, "done": 1, "total": 1}
            ],
            "story_points_alloc": [
                {"epic": "ADSI Data Ingestion", "points": 13, "status": "In Progress"},
                {"epic": "Data Streaming", "points": 8, "status": "In Progress"},
                {"epic": "Member Segmentation", "points": 5, "status": "To Do"},
                {"epic": "CHF Tag Rules", "points": 3, "status": "To Do"}
            ],
            "accomplishments": [
                "Built the ADSI data integration pipeline to bring targeting data into Data Cloud",
                "Implemented data streaming and mapping so ADSI fields populate Salesforce Data Cloud Objects",
                "Configured initial CHF tag-based segmentation rules"
            ],
            "risks": [
                {"title": "Targeting Logic & CTA Pending",
                 "description": "Finalization of targeting logic and Call-to-Action templates pending business review and approval",
                 "status": "Active Risk"}
            ],
            "next_steps": [
                {"title": "Approved Message Templates", "description": "Build the template library to support the pilot messaging campaigns"},
                {"title": "Enable SMS messaging Channel", "description": "Complete SMS channel setup and validate SMS verbiage and delivery"},
                {"title": "Activate Data Cloud → SFMC Mapping", "description": "Configure and validate Data Cloud and SFMC activation for outreach"}
            ]
        }
    ]
}

if __name__ == '__main__':
    buf = generate_sprint_deck(SAMPLE_DATA)
    with open('sprint_deck_test.pptx', 'wb') as f:
        f.write(buf.read())
    print("✅ Sprint deck generated: sprint_deck_test.pptx")